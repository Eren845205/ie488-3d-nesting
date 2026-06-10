"""make_pptx_3d.py — Aşama 2 (3D voxel nesting) sunumu üretir.

2D sunumun (make_pptx.py) görsel kalıbıyla aynı; içerik PLAN_3D.md +
results/summary_3d.md çıktılarından.  Konuşma metni her slaytın PowerPoint
"konuşmacı notları"na gömülüdür.

Çalıştır:  python scripts/make_pptx_3d.py
Çıktı   :  SUNUM_3D.pptx (proje kök dizini)

Gerekli görseller (önce üret):
  python scripts/generate_models.py
  python scripts/render_models_overview.py
  python -m src.nesting3d.run3d --scenario all --export-stl --iters 1000
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

_ROOT = Path(__file__).resolve().parent.parent
_RES = _ROOT / "results"

ACCENT = RGBColor(0x25, 0x63, 0xEB)
GREEN = RGBColor(0x05, 0x96, 0x59)
INK = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
BG = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEE, 0xF2, 0xFF)
WINBG = RGBColor(0xDC, 0xFC, 0xE7)
NOTEBG = RGBColor(0xFE, 0xF9, 0xC3)

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def _bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _accent_bar(slide):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SW, Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def _title(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = ACCENT
    ln = slide.shapes.add_shape(1, Inches(0.72), Inches(1.45), Inches(7.5), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()


def _bullets(slide, items, left=0.8, top=1.85, width=11.7, size=20):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, color, bold) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        r = p.add_run(); r.text = "•  " + text
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
    return tb


def _notebox(slide, text, left=0.8, top=5.5, width=11.7, height=1.2, size=16,
             bold=False):
    nb = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width),
                                Inches(height))
    nb.fill.solid(); nb.fill.fore_color.rgb = NOTEBG; nb.line.fill.background()
    ntf = nb.text_frame; ntf.word_wrap = True
    ntf.margin_left = Pt(12); ntf.margin_top = Pt(8)
    nr = ntf.paragraphs[0].add_run()
    nr.text = text
    nr.font.size = Pt(size); nr.font.color.rgb = INK; nr.font.bold = bold


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _pic_center(slide, path, top=1.7, width=10.0):
    if not path.exists():
        return
    left = Inches((13.333 - width) / 2)
    slide.shapes.add_picture(str(path), left, Inches(top), width=Inches(width))


def _caption(slide, text, left, top, width=5.8):
    cap = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width),
                                   Inches(0.5))
    cr = cap.text_frame.paragraphs[0].add_run()
    cr.text = text
    cr.font.size = Pt(12); cr.font.color.rgb = MUTED


# ---------------------------------------------------------------- 1. Başlık
s = prs.slides.add_slide(BLANK); _bg(s)
box = s.shapes.add_shape(1, Inches(0.7), Inches(1.0), Inches(4.8), Inches(0.55))
box.fill.solid(); box.fill.fore_color.rgb = LIGHT; box.line.fill.background()
btf = box.text_frame; btf.margin_top = Pt(2); btf.margin_bottom = Pt(2)
bp = btf.paragraphs[0]; br = bp.add_run()
br.text = "IE 488 — Dönem Projesi · Aşama 2"
br.font.size = Pt(16); br.font.bold = True; br.font.color.rgb = ACCENT
bp.alignment = PP_ALIGN.CENTER

tb = s.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(12), Inches(2.2))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; r = p.add_run()
r.text = "Eklemeli İmalatta 3B Voxel Nesting"
r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = INK

tb2 = s.shapes.add_textbox(Inches(0.7), Inches(4.2), Inches(12), Inches(1.8))
tf2 = tb2.text_frame; tf2.word_wrap = True
p2 = tf2.paragraphs[0]; r2 = p2.add_run()
r2.text = ("Karmaşık 3B modelleri voxel'leştirip yüksekliği minimize edilen "
           "kutuya yerleştirme — 2B iskeletin hacimsel genellemesi")
r2.font.size = Pt(20); r2.font.color.rgb = MUTED
p3 = tf2.add_paragraph(); p3.space_before = Pt(14)
r3 = p3.add_run()
r3.text = ("STL mesh  →  Voxelization  →  Order  →  DBLF + Simulated Annealing"
           "  →  .STL   ·   220×220 mm taban, yükseklik serbest")
r3.font.size = Pt(14); r3.font.color.rgb = MUTED
_notes(s,
"AÇILIŞ (3 cümle):\n"
"1) Geçen hafta 2B nesting'i göstermiştim: baseline sezgiseller + Simulated Annealing.\n"
"2) Hocamın yönlendirmesiyle bu hafta aynı iskeleti 3B'ye taşıdım: gerçek 3B modeller, "
"voxelization, yüksekliği minimize edilen kutu, çıktı olarak STL.\n"
"3) İki senaryoda test ettim: direktifteki sette baseline zaten çok güçlü çıktı; "
"daha kalabalık sette SA yüksekliği 3.4 mm düşürdü.\n"
"NOT: 'mükemmel/optimal çözdüm' deme — 'bu çözünürlükte optimuma yakın' de.")

# ------------------------------------------------- 2. Direktif (el yazısı not)
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Çıkış Noktası — Verilen Yön")
directive = _ROOT / "WhatsApp Image 2026-06-09 at 21.28.07.jpeg"
if directive.exists():
    s.shapes.add_picture(str(directive), Inches(0.8), Inches(1.75),
                         height=Inches(4.3))
_caption(s, "Toplantı notu (2026-06-09)", 0.8, 6.1)
tb = s.shapes.add_textbox(Inches(7.0), Inches(1.9), Inches(5.6), Inches(4.6))
tf = tb.text_frame; tf.word_wrap = True
for i, (txt, c, b) in enumerate([
    ("Okuduğum pipeline:", INK, True),
    ("1. 3B karmaşık form (3D mesh / STL)", INK, False),
    ("2. Voxelization — grid temsiline geçiş", INK, False),
    ("3. Order — sıralama + yerleştirme", INK, False),
    ("4. Çıktı: .STL", INK, False),
    ("Kutu: taban sabit, yükseklik serbest → yüksekliği minimize et", GREEN, True),
    ("3 model, adetler ×10 / ×5 / ×15", GREEN, True),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(10)
    r = p.add_run(); r.text = txt
    r.font.size = Pt(18); r.font.color.rgb = c; r.font.bold = b
_notes(s,
"'Geçen toplantıda hocam bu pipeline'ı çizdi: karmaşık bir 3B mesh alınacak, "
"voxel'leştirilecek, sıralanıp kutuya yerleştirilecek, sonuç STL olarak verilecek. "
"Kutunun tabanı sabit, yüksekliği serbest — amaç yüksekliği indirmek. "
"Üç model seçilecek, 10-5-15 adetlerle.'\n"
"Bu slayt 'direktifi doğru anladım mı' teyididir — hoca düzeltmek isterse burada düzeltir.")

# ---------------------------------------------------------------- 3. Problem
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Problem: 3B Strip Packing (Open-Dimension)")
_bullets(s, [
    ("Sabit 220×220 mm taban (2B fazla aynı tabla) — yükseklik sınırsız ama "
     "cezalı: amaç max yüksekliği minimize etmek.", INK, False),
    ("SLM/SLS toz yatağı mantığı: bina yüksekliği ≈ baskı süresi + toz maliyeti.", INK, False),
    ("2B'deki 'utilization maksimize' amacının open-dimension karşılığı.", INK, False),
    ("Literatür konumu: Tang vd. 2025'te N-IM (voxel) dalı — 2B'de E\\P\\S "
     "kutusundaydım, şimdi N\\P\\S'ye yükseldim.", INK, False),
    ("İkincil metrik: packing density = parça hacmi / kullanılan zarf hacmi.", INK, False),
])
_notes(s,
"'Problem 3B strip packing: taban sabit, yükseklik serbest, yüksekliği minimize ediyorum. "
"Toz yataklı sistemlerde yükseklik doğrudan baskı süresi ve toz maliyeti demek.'\n"
"HOCA SORARSA (2B'den farkı?): 2B'de alan doluluğunu maksimize ediyordum; "
"3B'de tavan açık olduğu için amaç yüksekliği indirmek — aynı ailenin open-dimension üyesi.\n"
"N-IM: review'de 'gerçek geometri temsili' sınıfı (voxel/polygon); "
"E-IM yalnız bounding box idi.")

# ---------------------------------------------------------------- 4. Modeller
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Parça Seti — 3 Model, 30 Parça")
_pic_center(s, _RES / "models_overview.png", top=1.8, width=11.5)
_notebox(s,
         "Modeller prosedürel üretiliyor (scripts/generate_models.py) — pipeline "
         "model-bağımsız: data/models/ klasörüne konan herhangi bir STL aynı akıştan geçer.",
         top=5.9, height=1.0)
_notes(s,
"'Üç model seçtim: hocamın çizdiği örneğe sadık kalarak bir sandalye (içbükey, "
"her yönden boşluklu), asimetrik bir L-braket ve ortası delik bir halka. "
"Adetler direktifteki gibi 10-5-15, toplam 30 parça.'\n"
"Neden bu üçü? Üçü de 'karmaşık form' sınıfında: içbükeylik, asimetri, delik — "
"bounding box'la temsil edilse büyük hacim israfı olurdu; voxel temsilin değerini gösteriyorlar.\n"
"Neden prosedürel? Tekrar üretilebilirlik + internete bağımlılık yok; "
"ama herhangi bir STL dosyası da sisteme verilebilir.")

# ---------------------------------------------------------------- 5. Yöntem 1
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Adım 1 — Voxelization")
_bullets(s, [
    ("Her mesh 3.44 mm'lik voxel'lere bölünür → taban 64×64 hücrelik grid "
     "(sandalye ≈ 11×11×21 voxel).", INK, False),
    ("4 eksen-hizalı duruş: dik / yan × 0° / 90° — her duruş için ayrı grid.", INK, False),
    ("Her duruşun kolon profilleri önceden çıkarılır: alt profil (bottom) ve "
     "üst profil (top) — yerleştirme bu profillerle O(taban) sürede.", INK, False),
    ("İç hacim dolduruluyor (fill) — yüzey-kabuk grid metrikleri bozardı.", INK, False),
    ("2B'deki raster/AABB fikrinin hacimsel hali: collision testi grid üzerinde.", INK, False),
])
_notebox(s,
         "Çözünürlük bilinçli kaba: 64×64 taban gridinde SA'nın her denemesi ~20 ms — "
         "1000 iterasyonluk arama dakikanın altında kalıyor.",
         top=5.7, height=1.1)
_notes(s,
"'Mesh'i 3.44 mm'lik voxel'lere bölüyorum; sandalye yaklaşık 11×11×21 voxel oluyor. "
"Her parçanın 4 duruşu var: dik ve yan, her biri 0 ve 90 derece. "
"Yerleştirme hızlı olsun diye her duruşun alt ve üst profillerini önceden çıkarıyorum.'\n"
"Neden 24 değil 4 oryantasyon? AM'de parça stabilitesi + arama uzayını küçük tutmak; "
"parametrik, artırılabilir.\n"
"Neden bu pitch? Hız/doğruluk dengesi: pitch küçülürse temsil hassaslaşır ama "
"SA iterasyonu pahalanır. 64'lük grid dakika-altı demoya izin veriyor.")

# ---------------------------------------------------------------- 6. Yöntem 2
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Adım 2 — DBLF Baseline + Heightmap")
_bullets(s, [
    ("Order: parçalar hacim-azalan sıralanır (2B BFD'nin alan-azalan kuralının "
     "3B'si) — direktifteki 'Order' adımı.", INK, False),
    ("Kutu durumu tek 2B yükseklik haritası: H[x,y] = o kolondaki en üst dolu voxel.", INK, False),
    ("Bir parça (x,y)'ye bırakılınca ineceği z, kolon profillerinden tek formülle "
     "bulunur — tüm pozisyonlar vektörize taranır.", INK, False),
    ("Seçim kuralı: tüm (duruş, x, y) adayları içinde en alçak TEPE: "
     "(z_tepe, z, y, x) sözlüksel minimumu = Deepest-Bottom-Left-Fill.", GREEN, True),
])
_notebox(s,
         "Bilinen sınırlama (bilinçli): heightmap, çıkıntı ALTINA parça sokamaz "
         "(örn. sandalye oturağının altı). Karşılığında her SA denemesi ucuz. "
         "Tam 3B voxel collision → Aşama 3 kapısı.",
         top=5.4, height=1.4)
_notes(s,
"'Kutunun durumunu tek bir yükseklik haritasıyla tutuyorum — Tetris gibi: "
"parça (x,y)'ye bırakılınca nereye oturacağı, kolon profillerinden tek formülle çıkıyor. "
"Baseline kuralım DBLF: bütün duruş ve pozisyon adayları içinde en alçak tepeyi "
"veren yerleşimi seç.'\n"
"Neden z değil z_tepe? Amaç tavanı indirmek: aynı z'de yatık duruş, dik duruştan iyi — "
"kural bunu direkt ödüllendiriyor.\n"
"SINIRLAMA SORULURSA: dürüst söyle — 'heightmap çıkıntı altı boşlukları kullanamaz; "
"bu, hız için bilinçli bir taviz. Tam 3B collision bir sonraki aşamanın işi.'")

# ---------------------------------------------------------------- 7. Yöntem 3 (SA)
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Adım 3 — Simulated Annealing Katmanı")
_bullets(s, [
    ("2B'deki SA iskeletinin portu — ama karar değişkeni büyüdü: "
     "sıra + her parçanın duruşu.", INK, False),
    ("Komşu hamleler: swap / insert / reverse (2B ile aynı) + duruş çevirme.", INK, False),
    ("Enerji = max yükseklik (mm) + 0.1 × RMS kolon yüksekliği.", GREEN, True),
    ("Başlangıç çözümü = DBLF'nin tam çözümü → SA asla baseline'ın altına düşemez.", INK, False),
    ("Deterministik: seed=42 ile her koşu aynı sonucu verir.", INK, False),
])
_notebox(s,
         "RMS neden? Ortalama yükseklik, dik duran ve yatan parçayı AYNI puanlar "
         "(aynı kolon-hacmi) — gradyan yok. Kareli ortalama yayvan/alçak yerleşimi "
         "kesin tercih eder; SA'ya plato dışına çıkacağı bir yokuş verir.",
         top=5.4, height=1.4)
_notes(s,
"'SA'nın 2B'deki iskeletini aynen taşıdım; tek fark karar değişkeni artık "
"sıra + duruş. Enerji fonksiyonum max yükseklik artı küçük bir kompaktlık terimi.'\n"
"EN GÜÇLÜ MÜHENDİSLİK HİKÂYESİ (sorulursa anlat): İlk denemede kompaktlık terimi "
"ortalama yükseklikti ve SA hiç ilerlemiyordu. Sebep: dik duran halka ile yatan halka "
"aynı kolon-hacmi kaplıyor — ortalamaya göre fark YOK, gradyan yok. "
"Kareli ortalamaya (RMS) geçince yatık duruş kesin avantajlı oldu ve arama çalıştı.\n"
"İkinci kalibrasyon: başlangıç sıcaklığı 8 mm'de kabul oranı %86'ydı — saf rastgele "
"yürüyüş. 3 mm'ye indirince arama dengelendi.")

# ------------------------------------------------- 8. Sonuç — default senaryo
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Sonuç 1 — Direktif Seti (30 parça)")
img = _RES / "nesting3d_default.png"
if img.exists():
    s.shapes.add_picture(str(img), Inches(0.55), Inches(1.7), width=Inches(7.4))
_caption(s, "DBLF = SA: 44.7 mm — halkalar dik, sandalyeler yatık", 0.8, 6.85)
tb = s.shapes.add_textbox(Inches(8.2), Inches(1.9), Inches(4.5), Inches(3.4))
tf = tb.text_frame; tf.word_wrap = True
for i, (txt, c, b) in enumerate([
    ("Yükseklik: 44.7 mm (13 voxel)", INK, True),
    ("SA 3000 iterasyonda dahi geçemedi → baseline bu çözünürlükte optimuma yakın.", INK, False),
    ("Greedy'nin bulduğu çözüm akıllıca: taban dolunca halkaları DİK dikip "
     "13 voxel'de tutuyor.", INK, False),
    ("Katı alt sınır: 11 voxel (yatan sandalyenin kalınlığı) — boşluk en çok 2 voxel.", INK, False),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(10)
    r = p.add_run(); r.text = ("•  " if i else "") + txt
    r.font.size = Pt(16); r.font.color.rgb = c; r.font.bold = b
_notebox(s, "2B'deki small set hikâyesinin 3B karşılığı: 'iyileşme yok' değil, "
            "'baseline doğrulandı'.", left=8.2, top=5.6, width=4.5, height=1.3,
         size=14, bold=True)
_notes(s,
"'Direktifteki sette baseline 44.7 mm buldu; SA'yı 3000 iterasyona kadar zorladım, geçemedi. "
"Bunu başarısızlık olarak değil doğrulama olarak okuyorum: greedy kuralın bulduğu çözüm "
"gerçekten akıllıca — tabanda yer kalmayınca halkaları dik dikiyor ve tavanı yükseltmiyor.'\n"
"ALT SINIR SORULURSA: 'Sette yatan sandalye var; sandalyenin en ince boyutu 11 voxel — "
"hiçbir yerleşim 11'in altına inemez. Ben 13'teyim; aradaki 2 voxel'in kapanabilir olup "
"olmadığını kolon-sayım argümanıyla daraltabiliyorum ama kanıtlamadım — dürüst boşluk bu.'\n"
"2B paralelliği kur: small set orada da 'zaten optimal' çıkmıştı; dürüst raporlama aynı çizgide.")

# -------------------------------------------------- 9. Sonuç — stress senaryo
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Sonuç 2 — Stress Seti (42 parça): SA Kazancı")
img = _RES / "nesting3d_stress.png"
if img.exists():
    s.shapes.add_picture(str(img), Inches(0.55), Inches(1.7), width=Inches(7.4))
_caption(s, "SA sonrası stress yerleşimi: 55.0 mm", 0.8, 6.85)
tb = s.shapes.add_textbox(Inches(8.2), Inches(1.9), Inches(4.5), Inches(4.4))
tf = tb.text_frame; tf.word_wrap = True
for i, (txt, c, b) in enumerate([
    ("DBLF: 58.4 mm", INK, False),
    ("SA: 55.0 mm (−3.4 mm)", GREEN, True),
    ("Density: 0.387 → 0.406", GREEN, True),
    ("Adetler ×1.4 (14/7/21): taban sıkışınca sıra + duruş kararları kritikleşiyor "
     "— SA'nın arama payı burada açılıyor.", INK, False),
    ("~25 saniyede 1000 iterasyon.", INK, False),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(10)
    r = p.add_run(); r.text = "•  " + txt
    r.font.size = Pt(16); r.font.color.rgb = c; r.font.bold = b
_notes(s,
"'Setin kalabalık versiyonunda — adetler 1.4 katı, 42 parça — tablo değişiyor: "
"baseline 58.4 mm'de kalıyor, SA bunu 55.0'a indiriyor. Bir voxel katmanı kazandı; "
"density de 0.39'dan 0.41'e çıktı.'\n"
"Neden burada kazanıyor? Taban sıkışınca greedy'nin erken kararları (kimi nereye, "
"hangi duruşla) sonraki parçaları kötü yerlere itiyor; SA sırayı ve duruşları birlikte "
"oynatarak bu tıkanıklığı çözüyor.\n"
"2B paralelliği: medium sette de durum buydu — kalabalık arttıkça arama değer üretiyor.")

# -------------------------------------------------- 10. Karşılaştırma grafiği
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Sonuçlar — DBLF vs SA Karşılaştırması")
_pic_center(s, _RES / "comparison_3d.png", top=1.85, width=11.6)
_notes(s,
"'Sol grafik yükseklik — alçak olan iyi: default'ta çubuklar eşit (baseline doğrulandı), "
"stress'te SA 3.4 mm indiriyor. Sağ grafik density — yüksek olan iyi: stress'te "
"0.387'den 0.406'ya çıkıyor.'\n"
"2B sunumdaki karşılaştırma grafiğinin 3B karşılığı — aynı anlatım: yeşil SA, "
"kalabalık sette baseline'ı geçiyor; seyrek sette eşit çünkü baseline zaten tavanda.\n"
"Eksen sorulursa: yükseklik mm cinsinden, density birimsiz oran (parça hacmi / "
"kullanılan zarf hacmi).")

# -------------------------------------------------- 11. Yakınsama + tablo
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Sayılarla — Yakınsama ve Özet Tablo")
img = _RES / "sa3d_convergence_stress.png"
if img.exists():
    s.shapes.add_picture(str(img), Inches(0.55), Inches(1.8), width=Inches(6.6))
_caption(s, "SA yakınsama (stress) — 90. iterasyonda iyileşme", 0.8, 5.6)
rows = [
    ["Senaryo", "Parça", "DBLF", "SA", "Kazanç"],
    ["default", "30", "44.7 mm", "44.7 mm", "— (baseline doğrulandı)"],
    ["stress", "42", "58.4 mm", "55.0 mm", "−3.4 mm · density +0.019"],
]
tbl = s.shapes.add_table(3, 5, Inches(7.45), Inches(2.1),
                         Inches(5.45), Inches(1.9)).table
widths = [1.25, 0.8, 1.05, 1.05, 1.30]
for c, w in enumerate(widths):
    tbl.columns[c].width = Inches(w)
for c in range(5):
    cell = tbl.cell(0, c)
    cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
    para = cell.text_frame.paragraphs[0]; run = para.add_run(); run.text = rows[0][c]
    run.font.size = Pt(13); run.font.bold = True; run.font.color.rgb = WHITE
    para.alignment = PP_ALIGN.CENTER
for r in range(1, 3):
    win = (r == 2)
    for c in range(5):
        cell = tbl.cell(r, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = WINBG if win else WHITE
        para = cell.text_frame.paragraphs[0]; run = para.add_run(); run.text = rows[r][c]
        run.font.size = Pt(12); run.font.bold = win
        run.font.color.rgb = RGBColor(0x16, 0x65, 0x34) if win else INK
        para.alignment = PP_ALIGN.CENTER
_notebox(s,
         "Dürüst okuma: kazanç tek hamlede geliyor (90. iterasyon) ve sonrası plato — "
         "voxel dünyasında iyileşmeler katman katmandır, sürekli eğri beklenmez.",
         left=7.45, top=4.3, width=5.45, height=1.6, size=14)
_notes(s,
"'Yakınsama grafiği SA'nın gerçekten arama yaptığının kanıtı: 90. iterasyon civarında "
"baseline'ın altına inen bir konfigürasyon buluyor ve onu koruyor.'\n"
"Neden tek basamak? Yükseklik voxel cinsinden ölçülüyor — 3.44 mm'lik katmanlar. "
"Bir sonraki iyileşme bir katman daha demek; sürekli, pürüzsüz bir eğri bu problemde olmaz.\n"
"Tekrar-üretilebilirlik: seed sabit, iki koşu aynı sonucu verir; tüm sayılar "
"results/summary_3d.md'de.")

# -------------------------------------------------- 12. STL çıktısı + demo
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Pipeline'ın Son Oku — .STL Çıktısı")
_bullets(s, [
    ("Yerleşim çözümü, voxel'lerden GERÇEK geometriye geri çevriliyor: her parçanın "
     "orijinal mesh'ine duruş + konum dönüşümü uygulanıp tek sahnede birleştiriliyor.", INK, False),
    ("results/nesting3d_result_default.stl · nesting3d_result_stress.stl", GREEN, True),
    ("Windows 3D Viewer / herhangi bir slicer ile açılabilir — baskıya giden "
     "dosya formatının kendisi.", INK, False),
    ("Tek tık demo: DEMO_3D.bat — modeller + iki senaryo + görseller + STL, ~1 dakika.", INK, False),
    ("Doğrulama: 72 birim testi (47 eski 2B + 25 yeni 3B) — 2B koduna dokunulmadı.", INK, False),
])
_notes(s,
"'Direktifteki son ok STL'di: çözümü voxel'de bırakmıyorum — her parçanın gerçek mesh'ine "
"bulunan duruş ve konumu uygulayıp tek STL sahnesi yazıyorum. Bu dosya doğrudan "
"slicer'a/baskıya gidebilir.' İsterseniz 3D Viewer'da açıp gösterebilirim.\n"
"Test vurgusu: 25 yeni birim testi var — yerleştirme çakışmazlığı, determinizm, "
"STL'in geri yüklenebilirliği, SA'nın baseline altına düşmemesi hepsi testle sabit. "
"2B kodu hiç değişmedi; 47 eski test hâlâ geçiyor.")

# -------------------------------------------------- 13. Sınırlar + sonraki adım
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Sınırlamalar ve Sonraki Adımlar")
_bullets(s, [
    ("Heightmap çıkıntı altını kullanamaz → tam 3B voxel collision ile 'cavity' "
     "doldurma (sandalye altına halka sokmak).", INK, False),
    ("Daha ince pitch (ör. 96×96 grid) — temsil hassaslaşır, süre/kalite eğrisi çıkarılır.", INK, False),
    ("Gerçek dünya STL'leri (lab parçaları) ile aynı pipeline'ın testi.", INK, False),
    ("Alt sınır analizi: 13 voxel'in optimal olduğunun kanıtı ya da 12'lik çözüm.", INK, False),
    ("2B'den taşınan diğer yön: GA karşılaştırması.", INK, False),
])
qb = s.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(11.7), Inches(1))
qr = qb.text_frame.paragraphs[0].add_run()
qr.text = "Hangisini önceliklendirmemi istersiniz?"
qr.font.size = Pt(24); qr.font.color.rgb = MUTED; qr.font.italic = True
_notes(s,
"'Beş yön görüyorum: cavity doldurma — yani sandalyenin altına halka sokabilmek; "
"daha ince çözünürlük; lab'ın gerçek parçalarıyla test; alt sınır analizi; ve GA karşılaştırması.'\n"
"KAPANIŞ HAMLESİ (mutlaka): hocaya dön ve sor — 'Hangisini önceliklendirmemi istersiniz?'\n"
"Neden cavity ilk sırada? En görünür sınırlama bu; çözülürse default sette de kazanç ihtimali doğar.")

# ---------------------------------------------------------------- 14. Özet
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Özet")
_bullets(s, [
    ("Direktifteki pipeline uçtan uca çalışıyor: STL → voxel → order → "
     "yükseklik-minimize kutu → STL.", GREEN, True),
    ("2B iskelet (constructive + SA) 3B'ye birebir taşındı — kod tabanı ayrık, "
     "2B'ye dokunulmadı.", INK, False),
    ("Direktif seti: 44.7 mm — baseline bu çözünürlükte optimuma yakın (SA doğruladı).", INK, False),
    ("Stress seti: 58.4 → 55.0 mm, SA'nın net arama kazancı.", INK, True),
    ("Tekrar-üretilebilir (seed=42) · 72 birim testi · tek tık demo + STL çıktıları.", INK, False),
])
qb = s.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(11.7), Inches(1))
qr = qb.text_frame.paragraphs[0].add_run()
qr.text = "Teşekkürler — sorularınız?"
qr.font.size = Pt(24); qr.font.color.rgb = MUTED
_notes(s,
"Kapanış (4 cümle): 'Çizdiğiniz pipeline'ı uçtan uca kurdum. 2B'de geliştirdiğim iskelet "
"3B'ye taşındı. Sizin verdiğiniz sette baseline çok güçlü çıktı ve SA bunu doğruladı; "
"kalabalık sette SA 3.4 mm kazandırdı. Hepsi tekrar-üretilebilir ve 72 testle doğrulanmış "
"durumda — STL çıktıları da hazır.'\n"
"Sonra teşekkür + soru al.")

out = _ROOT / "SUNUM_3D.pptx"
prs.save(str(out))
print(f"Yazildi: {out}  ({len(prs.slides)} slayt)")
