"""make_pptx.py — SUNUM.html içeriğinden PowerPoint (SUNUM.pptx) üretir.

Telefondan / PowerPoint'ten sunum için. HTML'deki konuşma notları, her slaytın
PowerPoint "konuşmacı notları"na gömülür.

Çalıştır:  python scripts/make_pptx.py
Çıktı   :  SUNUM.pptx (proje kök dizini)
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

_ROOT = Path(__file__).resolve().parent.parent
_RES = _ROOT / "results"

ACCENT = RGBColor(0x25, 0x63, 0xEB)
GREEN = RGBColor(0x05, 0x96, 0x59)
INK = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
BG = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEE, 0xF2, 0xFF)
WINBG = RGBColor(0xDC, 0xFC, 0xE7)
NOTEBG = RGBColor(0xFE, 0xF9, 0xC3)

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def _bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _accent_bar(slide):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SW, Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def _title(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = ACCENT
    ln = slide.shapes.add_shape(1, Inches(0.72), Inches(1.45), Inches(7.5), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()


def _bullets(slide, items, left=0.8, top=1.85, width=11.7, size=20):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, color, bold) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        r = p.add_run(); r.text = "•  " + text
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
    return tb


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _pic_center(slide, path, top=1.7, width=10.0):
    if not path.exists():
        return
    w = Inches(width)
    left = Emu(int((SW - w) / 1))
    left = Inches((13.333 - width) / 2)
    slide.shapes.add_picture(str(path), left, Inches(top), width=w)


# ---------------------------------------------------------------- 1. Başlık
s = prs.slides.add_slide(BLANK); _bg(s)
box = s.shapes.add_shape(1, Inches(0.7), Inches(1.0), Inches(4.2), Inches(0.55))
box.fill.solid(); box.fill.fore_color.rgb = LIGHT; box.line.fill.background()
btf = box.text_frame; btf.margin_top = Pt(2); btf.margin_bottom = Pt(2)
bp = btf.paragraphs[0]; br = bp.add_run(); br.text = "IE 488 — Dönem Projesi"
br.font.size = Pt(16); br.font.bold = True; br.font.color.rgb = ACCENT
bp.alignment = PP_ALIGN.CENTER

tb = s.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(12), Inches(2.2))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; r = p.add_run()
r.text = "Eklemeli İmalatta 2B Nesting Optimizasyonu"
r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = INK

tb2 = s.shapes.add_textbox(Inches(0.7), Inches(4.2), Inches(12), Inches(1.5))
tf2 = tb2.text_frame; tf2.word_wrap = True
p2 = tf2.paragraphs[0]; r2 = p2.add_run()
r2.text = ("Parçaları build tablasına en az boşlukla yerleştirme: "
           "baseline sezgisellerden metasezgisel optimizasyona")
r2.font.size = Pt(20); r2.font.color.rgb = MUTED
p3 = tf2.add_paragraph(); p3.space_before = Pt(14)
r3 = p3.add_run()
r3.text = "Constructive heuristic (BL + BFD)  →  Simulated Annealing  ·  220×220 mm FDM tablası"
r3.font.size = Pt(14); r3.font.color.rgb = MUTED
_notes(s,
"AÇILIŞ (3 cümle):\n"
"1) AM'de nesting problemine baktım: parçaları tablaya en az boşlukla dizmek — bir Endüstri Mühendisliği optimizasyon problemi.\n"
"2) Önce iki temel yöntem yazdım (Bottom-Left + Best-Fit Decreasing), sonra üstüne bir Simulated Annealing optimizasyon katmanı geliştirdim.\n"
"3) Medium sette doluluğu %86.9'dan %90.9'a çıkardım.\n"
"NOT: 'en optimize' deme — 'doluluğu artıracak şekilde' de.")

# ---------------------------------------------------------------- 2. Problem
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s); _title(s, "Problem: Nesting")
_bullets(s, [
    ("Bir 3D yazıcı tablasına birden çok parçayı yerleştir.", INK, False),
    ("Amaç: alan kullanımını maksimize et → atığı, malzemeyi, baskı süresini azalt.", INK, False),
    ("Klasik bir Endüstri Mühendisliği optimizasyon problemi: kısıtlı kaynağın (tabla) en verimli kullanımı.", INK, False),
    ("Kapsam: NfAM (yalnız yerleştirme) · 2B · tek tabla.", INK, False),
])
nb = s.shapes.add_shape(1, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.0))
nb.fill.solid(); nb.fill.fore_color.rgb = NOTEBG; nb.line.fill.background()
ntf = nb.text_frame; ntf.word_wrap = True; ntf.margin_left = Pt(12)
np_ = ntf.paragraphs[0]; nr = np_.add_run()
nr.text = "Benzetme: valize en az boşlukla eşya yerleştirmek ya da otoparka en çok arabayı sığdırmak."
nr.font.size = Pt(16); nr.font.color.rgb = INK
_notes(s,
"'Nesting = parçaları tablaya dizme. Amaç en az boşluk, yani en az atık ve süre. Kısıtlı kaynağı (tabla) verimli kullanma — tam bir IE konusu.'\n"
"Benzetme ver: valize en az boşlukla eşya koymak.\n"
"HOCA SORARSA (Neden önemli?): Az atık = az malzeme + az baskı süresi + az maliyet.")

# ---------------------------------------------------------------- 3. Kapsam
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s); _title(s, "Kapsam ve Yaklaşım")
_bullets(s, [
    ("Parça temsili: AABB — her parça bir dikdörtgen (en × boy). Rotation {0°, 90°}.", INK, False),
    ("Tabla: 220×220 mm FDM yatağı, parçalar arası 1 mm boşluk.", INK, False),
    ("Literatür konumu: Tang vd. 2025 review'inin E\\P\\S kutusu (Extended mapping / Pure placement / Single).", INK, False),
    ("Araujo 2019 taksonomisi de bu sınıf için Bottom-Left tipi sezgiseli öneriyor → seçimim literatürle uyumlu.", INK, False),
])
_notes(s,
"'Parçayı dikdörtgenle temsil ettim (AABB) — en yaygın, en ucuz. 2B seçtim çünkü FDM'de izdüşüm yeterli ve aynı iskelet 3D'ye taşınabilir.'\n"
"E\\P\\S: E=dikdörtgen temsil (AABB), P=sadece yerleştir, S=tek algoritma.\n"
"Neden AABB? En basit, en hızlı, 'en eski en yaygın'. Polygon sonraki adım.")

# ---------------------------------------------------------------- 4. Baseline
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Adım 1 — Baseline (Constructive Sezgiseller)")
tb = s.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(6.2), Inches(5))
tf = tb.text_frame; tf.word_wrap = True
for i, (h, body) in enumerate([
    ("Bottom-Left (BL)", "Parçaları sol-alttan yerleştirir; en temel referans."),
    ("Best-Fit Decreasing (BFD)", "Alana göre sıralar, en az artık bırakan yere koyar (skyline)."),
]):
    ph = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    ph.space_before = Pt(6 if i else 0)
    rh = ph.add_run(); rh.text = h; rh.font.size = Pt(20); rh.font.bold = True; rh.font.color.rgb = GREEN
    pb = tf.add_paragraph(); pb.space_after = Pt(10)
    rb = pb.add_run(); rb.text = body; rb.font.size = Pt(18); rb.font.color.rgb = INK
nb = s.shapes.add_shape(1, Inches(0.8), Inches(5.4), Inches(6.2), Inches(1.2))
nb.fill.solid(); nb.fill.fore_color.rgb = NOTEBG; nb.line.fill.background()
ntf = nb.text_frame; ntf.word_wrap = True; ntf.margin_left = Pt(10)
nr = ntf.paragraphs[0].add_run()
nr.text = "İkisi de açgözlü: anlık iyi olanı seçer, en iyiyi garanti etmez."
nr.font.size = Pt(16); nr.font.color.rgb = INK; nr.font.bold = True
img = _RES / "bl_medium_on.png"
if img.exists():
    s.shapes.add_picture(str(img), Inches(7.35), Inches(1.75), height=Inches(5.0))
cap = s.shapes.add_textbox(Inches(7.35), Inches(6.8), Inches(5.6), Inches(0.5))
cr = cap.text_frame.paragraphs[0].add_run()
cr.text = "Baseline (BL) — medium set, %86.9 · gri taralı alan = atık"
cr.font.size = Pt(12); cr.font.color.rgb = MUTED
_notes(s,
"'İki temel yöntem yazdım. Bottom-Left: sol-alttan yerleştirir. Best-Fit Decreasing: alana göre sıralayıp en az artık bırakan yere koyar. İkisi de açgözlü — en iyiyi garanti etmez. Bu sınırı aşmak için optimizasyon ekledim.'\n"
"Görseli göster: gri taralı = atık.\n"
"Greedy ne? Her adımda o an en iyiyi seçen, bütünü düşünmeyen yöntem.")

# ---------------------------------------------------------------- 5. Optimizasyon
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Adım 2 — Geliştirdiğim Optimizasyon")
_bullets(s, [
    ("Gözlem: açgözlü kural, parçaların hangi sırada verildiğine duyarlı.", INK, False),
    ("O hâlde parça sırasını bir karar değişkeni yapıp en iyi sırayı arayabilirim.", GREEN, True),
    ("Yöntem: Simulated Annealing (SA) — içeride yine Bottom-Left'i kullanır, baseline'a dokunulmadı.", INK, False),
])
nb = s.shapes.add_shape(1, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.5))
nb.fill.solid(); nb.fill.fore_color.rgb = NOTEBG; nb.line.fill.background()
ntf = nb.text_frame; ntf.word_wrap = True; ntf.margin_left = Pt(12); ntf.margin_top = Pt(8)
nr = ntf.paragraphs[0].add_run()
nr.text = ("SA nedir? Farklı parça sıralarını deneyip en çok sığdıran sırayı arar; "
           "ara sıra kötü bir sırayı da kabul eder ki tek bir yere takılıp kalmasın.")
nr.font.size = Pt(17); nr.font.color.rgb = INK
_notes(s,
"EN KRİTİK SLAYT.\n"
"'Fark ettim ki açgözlü kural parça sırasına duyarlı. O yüzden sırayı karar değişkeni yapıp en iyi sırayı aradım — Simulated Annealing ile.'\n"
"TEK-NEFESLİK TANIM (ezberle): 'Farklı parça sıralarını deneyip en çok sığdıranı arar; ara sıra kötü sırayı da kabul eder ki takılıp kalmasın.'\n"
"DERİN SORARSA (hamle/soğutma/kabul): 'Bu yönü geliştiriyorum; mantığı: sırayı değiştir, doluluğu ölç, daha iyiyi tut, ara sıra kötüyü de kabul et.' — bu kadar yeter.")

# ---------------------------------------------------------------- 6. Karşılaştırma grafiği
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Sonuçlar — Doluluk Karşılaştırması")
_pic_center(s, _RES / "comparison_with_sa.png", top=1.75, width=10.5)
_notes(s,
"'Yeşil çubuklar SA. Medium ve stress'te baseline'ı geçiyor; small'da eşit çünkü orada tüm parçalar zaten sığıyor — baseline optimal.'\n"
"x ekseni = parça setleri, y ekseni = doluluk %.")

# ---------------------------------------------------------------- 7. Tablo
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Sonuçlar — Sayılarla")
rows = [
    ["Set", "Talep", "Baseline (en iyi)", "SA", "Yorum"],
    ["small", "%48", "47.6%", "47.6%", "Zaten optimal (hepsi sığıyor)"],
    ["medium", "%111", "86.9%", "90.9%", "Aramanın net kazancı: +4 puan"],
    ["stress", "%181", "82.7%", "91.7%", "%91.7'ye çıktı (çoğu sıralamadan)"],
]
tbl = s.shapes.add_table(4, 5, Inches(0.65), Inches(1.85),
                         Inches(12.0), Inches(2.6)).table
widths = [1.6, 1.4, 2.8, 1.4, 4.8]
for c, w in enumerate(widths):
    tbl.columns[c].width = Inches(w)
for c in range(5):
    cell = tbl.cell(0, c)
    cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
    para = cell.text_frame.paragraphs[0]; run = para.add_run(); run.text = rows[0][c]
    run.font.size = Pt(16); run.font.bold = True; run.font.color.rgb = WHITE
    para.alignment = PP_ALIGN.CENTER
for r in range(1, 4):
    win = (r == 2)
    for c in range(5):
        cell = tbl.cell(r, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = WINBG if win else (LIGHT if r % 2 == 0 else WHITE)
        para = cell.text_frame.paragraphs[0]; run = para.add_run(); run.text = rows[r][c]
        run.font.size = Pt(15); run.font.bold = win
        run.font.color.rgb = RGBColor(0x16, 0x65, 0x34) if win else INK
        para.alignment = PP_ALIGN.CENTER
nb = s.shapes.add_shape(1, Inches(0.65), Inches(4.9), Inches(12.0), Inches(1.5))
nb.fill.solid(); nb.fill.fore_color.rgb = NOTEBG; nb.line.fill.background()
ntf = nb.text_frame; ntf.word_wrap = True; ntf.margin_left = Pt(12); ntf.margin_top = Pt(8)
nr = ntf.paragraphs[0].add_run()
nr.text = ("Dürüst okuma: En temiz örnek medium — SA, baseline'ın bile altından başlayıp +4 puana "
           "tırmanıyor. Stress'te yüksek doluluğa ulaştım ama kazancın çoğu doğru sıralamadan, bir kısmı aramadan.")
nr.font.size = Pt(15); nr.font.color.rgb = INK
_notes(s,
"'Talep = tüm parçaların alanı / tabla. %100'ü aşınca bir kısmı zaten sığmaz. small %48 hepsi sığar. medium %111 SA +4 puan, en temiz örnek. stress %181 → %91.7 ama dürüst: kazancın çoğu sıralamadan.'\n"
"Doluluk %86 ama parça yerleşmedi? Doluluk=yerleşen alan/tabla; talep %111 olduğu için kısmı zaten sığmaz.\n"
"+9 hep SA'dan mı? Hayır — stress ~+7 sıralamadan, ~+2 aramadan. Temiz örnek medium.")

# ---------------------------------------------------------------- 8. Yakınsama
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "SA Nasıl İyileştiriyor? (Yakınsama)")
_pic_center(s, _RES / "sa_convergence.png", top=1.75, width=9.6)
_notes(s,
"'Bu grafik SA'nın çalıştığının kanıtı: iterasyonlar boyunca en iyi doluluk tırmanıyor. Medium (mavi) baseline'ın altından başlayıp onu geçiyor — arama gerçekten daha iyi sıra buluyor. small (yeşil) düz çünkü zaten tavanda.'")

# ---------------------------------------------------------------- 9. Sonraki adımlar
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Sonraki Adımlar")
_bullets(s, [
    ("(a) Genetik Algoritma — aynı mantıkla popülasyon tabanlı arama.", INK, False),
    ("(b) Polygon temsil + No-Fit-Polygon — dikdörtgen yerine gerçek kontur.", INK, False),
    ("(c) 3D voxel temsile geçiş — gerçek hacimsel nesting (SLM toz yatağı).", INK, False),
])
qb = s.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11.7), Inches(1))
qr = qb.text_frame.paragraphs[0].add_run()
qr.text = "Hangisini önceliklendirmemi istersiniz?"
qr.font.size = Pt(24); qr.font.color.rgb = MUTED; qr.font.italic = True
_notes(s,
"'Üç yöne gidebilirim.' (a/b/c'yi tek cümleyle söyle.) Sonra hocaya DÖN ve SOR: 'Hangisini önceliklendirmemi istersiniz?' — kapanış hamlesi, mutlaka sor.\n"
"Neden şimdi yapmadın? 'Önce yönü sizinle netleştirmek istedim ki doğru şeye emek harcayayım.'")

# ---------------------------------------------------------------- 10. Özet
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Özet")
_bullets(s, [
    ("Çalışan bir 2B AM nesting baseline'ı kurdum (BL + BFD, E\\P\\S kutusu).", INK, False),
    ("Üstüne bir Simulated Annealing optimizasyon katmanı geliştirdim.", GREEN, True),
    ("Kısıtlı sette doluluk: 86.9% → 90.9%  (+4 puan, medium).", INK, True),
    ("Tekrar-üretilebilir (seed=42), 47 birim testiyle doğrulandı.", INK, False),
])
qb = s.shapes.add_textbox(Inches(0.8), Inches(5.4), Inches(11.7), Inches(1))
qr = qb.text_frame.paragraphs[0].add_run()
qr.text = "Teşekkürler — sorularınız?"
qr.font.size = Pt(24); qr.font.color.rgb = MUTED
_notes(s,
"Kapanış: baseline kurdum → üstüne SA geliştirdim → medium'da +4 puan, tekrar-üretilebilir + 47 testle doğrulandı. 'Teşekkürler, sorularınız?'")

out = _ROOT / "SUNUM.pptx"
prs.save(str(out))
print(f"Yazildi: {out}  ({len(prs.slides)} slayt)")
