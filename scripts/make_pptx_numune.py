"""make_pptx_numune.py — Numune fazı (gerçek parçalar) sonuç sunumu üretir.

make_pptx_3d.py görsel kalıbıyla aynı; içerik hocanın yolladığı 8 gerçek
numunenin (48 parça) FİNAL yerleştirme sonuçları (2026-06-12 — hoca feedback
şartları uygulanmış: 335×335×600 makine, >=1 mm parça arası boşluk,
devoxelization ile orijinal çözünürlük).  Konuşma metni her slaytın
PowerPoint "konuşmacı notları"na gömülüdür.

Çalıştır:  python scripts/make_pptx_numune.py
Çıktı   :  SUNUM_NUMUNE.pptx (proje kök dizini)

Gerekli görseller (önce üret):
  python -m src.nesting3d.run3d --scenario numune --algo dblf --pitch 1.5 \
      --orient hybrid --export-stl --check-clearance --out results/numune_dblf
  python -m src.nesting3d.run3d --scenario numune --algo sa --iters 2000 \
      --pitch 1.5 --orient hybrid --seed 13 --export-stl --check-clearance \
      --out results/numune_sa
  python scripts/compare_numune.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

_ROOT = Path(__file__).resolve().parent.parent
_RES = _ROOT / "results"

ACCENT = RGBColor(0x25, 0x63, 0xEB)
GREEN = RGBColor(0x05, 0x96, 0x59)
DGREEN = RGBColor(0x16, 0x65, 0x34)
RED = RGBColor(0xB9, 0x1C, 0x1C)
INK = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
BG = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEE, 0xF2, 0xFF)
WINBG = RGBColor(0xDC, 0xFC, 0xE7)
NOTEBG = RGBColor(0xFE, 0xF9, 0xC3)

# --- FİNAL koşu rakamları (results/numune_*/summary_3d.csv, 2026-06-12) ---
PARTS = 48
PLATE = "335×335 mm"          # hoca makinesi tabanı
MAX_Z = "600 mm"              # hoca makinesi z limiti
PITCH = "1.5 mm"
MARGIN = ">= 1 mm (garantili)"
DBLF_H, DBLF_D = "214.5 mm", "0.189"
SA_H, SA_D = "181.5 mm", "0.229"
GAIN = "−33.0 mm  (%15.4)"
SA_TIME = "~93 dk (2000 iterasyon, seed 13)"
CLEAR_MM = "1.50 mm"          # clearance.py ölçümü, orijinal mesh'lerde
CLEAR_PAIRS = "46 çift"
TRI = "3.950.364"             # 48 parçanın toplam üçgeni = çıktı STL üçgeni
HOCA_HEDEF = "170 mm"         # sözlü hedef — ulaşılamadı, dürüst çerçevede

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def _bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _accent_bar(slide):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SW, Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def _title(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = ACCENT
    ln = slide.shapes.add_shape(1, Inches(0.72), Inches(1.45), Inches(7.5), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()


def _bullets(slide, items, left=0.8, top=1.85, width=11.7, size=20):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, color, bold) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        r = p.add_run(); r.text = "•  " + text
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
    return tb


def _notebox(slide, text, left=0.8, top=5.5, width=11.7, height=1.2, size=16,
             bold=False):
    nb = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width),
                                Inches(height))
    nb.fill.solid(); nb.fill.fore_color.rgb = NOTEBG; nb.line.fill.background()
    ntf = nb.text_frame; ntf.word_wrap = True
    ntf.margin_left = Pt(12); ntf.margin_top = Pt(8)
    nr = ntf.paragraphs[0].add_run()
    nr.text = text
    nr.font.size = Pt(size); nr.font.color.rgb = INK; nr.font.bold = bold


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _pic_center(slide, path, top=1.7, width=10.0):
    if not path.exists():
        return
    left = Inches((13.333 - width) / 2)
    slide.shapes.add_picture(str(path), left, Inches(top), width=Inches(width))


def _pic_center_h(slide, path, top, height):
    """Yükseklik-sınırlı ortalanmış görsel — slayt/notebox taşmasını önler."""
    if not path.exists():
        return
    from PIL import Image
    w, h = Image.open(path).size
    width = height * w / h
    left = Inches((13.333 - width) / 2)
    slide.shapes.add_picture(str(path), left, Inches(top), height=Inches(height))


def _caption(slide, text, left, top, width=6.5):
    cap = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width),
                                   Inches(0.5))
    cr = cap.text_frame.paragraphs[0].add_run()
    cr.text = text
    cr.font.size = Pt(12); cr.font.color.rgb = MUTED


def _metric_panel(slide, rows, left=8.2, top=1.9, width=4.6):
    """Sağ panel: büyük, net rakamlar (etiket + değer çiftleri)."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width),
                                  Inches(4.2))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for label, value, color, big in rows:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(2)
        r = p.add_run(); r.text = label
        r.font.size = Pt(15); r.font.color.rgb = MUTED
        p2 = tf.add_paragraph(); p2.space_after = Pt(14)
        r2 = p2.add_run(); r2.text = value
        r2.font.size = Pt(34 if big else 24); r2.font.bold = True
        r2.font.color.rgb = color
    return tb


# ---------------------------------------------------------------- 1. Başlık
s = prs.slides.add_slide(BLANK); _bg(s)
box = s.shapes.add_shape(1, Inches(0.7), Inches(0.9), Inches(6.4), Inches(0.55))
box.fill.solid(); box.fill.fore_color.rgb = LIGHT; box.line.fill.background()
btf = box.text_frame; btf.margin_top = Pt(2); btf.margin_bottom = Pt(2)
bp = btf.paragraphs[0]; br = bp.add_run()
br.text = "IE 488 — Dönem Projesi · Numune Fazı FİNAL (Makine Şartlarıyla)"
br.font.size = Pt(16); br.font.bold = True; br.font.color.rgb = ACCENT
bp.alignment = PP_ALIGN.CENTER

tb = s.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(12), Inches(1.6))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; r = p.add_run()
r.text = "Gerçek Numunelerle 3B Nesting — Final Sonuç Raporu"
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = INK

# büyük sonuç bandı
band = s.shapes.add_shape(1, Inches(0.7), Inches(3.3), Inches(12), Inches(1.7))
band.fill.solid(); band.fill.fore_color.rgb = WINBG; band.line.fill.background()
btf2 = band.text_frame; btf2.word_wrap = True
btf2.margin_left = Pt(20); btf2.margin_top = Pt(10)
bp1 = btf2.paragraphs[0]
br1 = bp1.add_run(); br1.text = f"Z yüksekliği:  {DBLF_H}  →  {SA_H}"
br1.font.size = Pt(34); br1.font.bold = True
br1.font.color.rgb = DGREEN
bp2 = btf2.add_paragraph()
br2 = bp2.add_run()
br2.text = (f"parça arası ölçülen min boşluk {CLEAR_MM} ✓   ·   "
            f"density {DBLF_D} → {SA_D}   ·   kabartma yazılar korunuyor "
            f"({TRI} üçgen)")
br2.font.size = Pt(18); br2.font.color.rgb = INK

tb2 = s.shapes.add_textbox(Inches(0.7), Inches(5.4), Inches(12), Inches(1.4))
tf2 = tb2.text_frame; tf2.word_wrap = True
p2 = tf2.paragraphs[0]; r2 = p2.add_run()
r2.text = (f"8 gerçek numune · toplam {PARTS} parça · makine {PLATE}, "
           f"z limiti {MAX_Z} · voxel {PITCH} · boşluk {MARGIN}")
r2.font.size = Pt(18); r2.font.color.rgb = MUTED
p3 = tf2.add_paragraph(); p3.space_before = Pt(8)
r3 = p3.add_run()
r3.text = ("STL → konservatif voxelization → DBLF → Simulated Annealing → "
           "devoxelization → orijinal çözünürlüklü .STL")
r3.font.size = Pt(14); r3.font.color.rgb = MUTED
_notes(s,
"AÇILIŞ: 'Geçen haftaki şartlarınızın üçünü de uyguladım: makine ölçüleri "
"335'e 335, 600 yükseklik; parça arası en az 1 milimetre boşluk — ölçülen 1.50; "
"ve çıktı orijinal çözünürlükte, kabartma yazılar duruyor. Bu şartlar altında "
"48 parça 181.5 milimetreye yerleşti.'\n"
"Önceki 155 rakamı sorulursa: o koşu 350 tabandaydı, parça arası boşluk yoktu ve "
"sonradan ölçtük — 0.77 mm'lik ihlal vardı. Yeni kurallarla geçerli değil; "
"şimdiki 181.5 ölçülmüş 1.50 mm boşlukla geçerli sonuç.")

# ------------------------------------------------- 2. Hoca şartları → durum
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Geçen Haftaki Şartlar → Durum")
rows = [
    ["Şart (2026-06-11)", "Uygulanan çözüm", "Durum"],
    ["Parça arası min 1 mm boşluk (her yönde)",
     f"Konservatif voxel sarması + yatay/dikey boşluk şeması; orijinal "
     f"geometride ölçüm: min {CLEAR_MM}", "✓ ölçüldü"],
    [f"Makine: taban 335×335 mm, max z {MAX_Z}",
     f"Taban {PLATE}; sonuç {SA_H} ≤ {MAX_Z}", "✓"],
    ["Çözünürlük: kabartma yazılar kaybolmayacak",
     f"Devoxelization — çıktı STL orijinal mesh'lerden; {TRI} üçgenin "
     f"tamamı birebir korunuyor", "✓ doğrulandı"],
    ["Optimizasyon voxel'de, çıktı orijinal parçalarla",
     "Voxel yalnız yerleşim hesabında; STL sahnesi orijinal geometriyle "
     "kurulur", "✓"],
]
tbl = s.shapes.add_table(5, 3, Inches(0.8), Inches(1.85),
                         Inches(11.7), Inches(4.0)).table
widths = [4.1, 6.0, 1.6]
for c, w in enumerate(widths):
    tbl.columns[c].width = Inches(w)
for c in range(3):
    cell = tbl.cell(0, c)
    cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
    para = cell.text_frame.paragraphs[0]; run = para.add_run()
    run.text = rows[0][c]
    run.font.size = Pt(15); run.font.bold = True; run.font.color.rgb = WHITE
for ri in range(1, 5):
    for c in range(3):
        cell = tbl.cell(ri, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
        para = cell.text_frame.paragraphs[0]; run = para.add_run()
        run.text = rows[ri][c]
        run.font.size = Pt(13)
        run.font.bold = (c == 2)
        run.font.color.rgb = DGREEN if c == 2 else INK
_notes(s,
"'Bu tablo geçen haftaki konuşmamızın birebir karşılığı. Dört maddenin dördü de "
"kapandı ve ikisi sadece kod değil ölçümle kanıtlı: boşluk yerleşim sonrası gerçek "
"geometri üzerinde ölçülüyor, üçgen sayısı çıktı dosyasında sayılıyor.'\n"
"Bu slayt sunumun omurgası — hoca ne istediyse onun karşılığını görüyor.")

# ------------------------------------------------- 3. Numune seti (görsel)
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Verilen Numuneler — 8 Model, 48 Parça")
_pic_center_h(s, _RES / "numune_overview.png", top=1.6, height=4.55)
_notebox(s,
         "Gerçek mm boyutları korunarak yüklendi (ölçekleme YOK). Taban artık "
         "makinenizin ölçüsü: 335×335 mm. Kritik geometri: 12 büyük delikli plaka "
         "(3/6/7/8 ×3) — taban aynı anda tek plaka alıyor, istif stratejisi "
         "sonucu belirliyor. Adetler: 1×4 · 2×16 · 3×3 · 4×15 · 5×1 · 6×3 · 7×3 · 8×3.",
         top=6.3, height=1.0, size=13)
_notes(s,
"'Sekiz numune gerçek boyutlarıyla. Belirleyici olan on iki büyük plaka: "
"335 milimetrelik tabana yan yana iki plaka sığmıyor — 228 artı 180, 408 eder. "
"Yani plakalar ister istemez üst üste istifleniyor ve yüksekliği o istif belirliyor.'\n"
"SORULURSA (plaka kalınlıkları): ham kalınlık toplamı 202.5 mm — buna boşluklar "
"eklenince salt yatık istif ~206+ mm tabana çarpıyor; 181.5'e inişin yolu bir "
"sonraki slaytlarda.")

# ------------------------------------------------- 4. Kurulum / yöntem özeti
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Kurulum — Makine Şartlarıyla Final Konfigürasyon")
_bullets(s, [
    (f"Taban {PLATE}, voxel {PITCH}, z limiti {MAX_Z} — makinenizin ölçüleri.",
     INK, False),
    ("Boşluk garantisi: konservatif voxelization (yüzey örnekleme — ince "
     "duvar/detay kaçmaz) + 1 voxel yatay genişletme + dikey istif boşluğu. "
     "Ardından bağımsız ölçüm: orijinal mesh'lerde en yakın çift araması.",
     GREEN, True),
    ("Poz stratejisi (hibrit): kısa plakalar (3, 6) dik durabilir; uzun "
     "plakalar (7, 8) yalnız yatık — dik boyları 187.5/190.1 mm tavan "
     "yaratıyordu. Bu seçim algoritmaya değil, parça analizine dayanıyor.",
     INK, False),
    ("Arama: DBLF baseline + SA (sıra + duruş, 2000 iterasyon) — "
     "best-so-far korunur, baseline'ın altına düşmez.", INK, False),
    ("Doğrulama: 95 birim testi yeşil; seed sabit, tekrar-üretilebilir.",
     INK, False),
])
_notebox(s,
         f"Raporlanan yükseklikler {PITCH} katmanlara YUKARI yuvarlanmış muhafazakâr "
         "değerlerdir. Boşluk rakamı ise voxel değil, gerçek geometri ölçümüdür.",
         top=5.95, height=1.0)
_notes(s,
"'Boşluk şartını iki katmanda çözdüm: önce voxel modelde garanti — parça gövdesi "
"bir hücre genişletiliyor, üst üste binmek matematiksel olarak imkânsız. Sonra "
"kontrol — yerleşim bittikten sonra gerçek geometri üzerinde en yakın iki parçanın "
"mesafesi ölçülüyor: 1.50 milimetre.'\n"
"Hibrit poz kararı sorulursa: plakaların dik boyları 178.3 / 180.9 / 187.5 / 190.1 — "
"uzunlara dik izni verilirse tavan 190'a kilitleniyor; kısalara izin verip uzunları "
"yatırınca 181.5 mümkün oldu.")

# ------------------------------------------------- 5. DBLF yerleşimi
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Yerleşim 1 — DBLF Baseline")
img = _RES / "numune_dblf" / "nesting3d_numune.png"
if img.exists():
    s.shapes.add_picture(str(img), Inches(0.45), Inches(1.65), width=Inches(7.5))
_caption(s, "48 parçanın tamamı yerleşti — greedy kural, hacim-azalan sıra", 0.8, 6.9)
_metric_panel(s, [
    ("Z yüksekliği", DBLF_H, INK, True),
    ("Packing density", DBLF_D, INK, True),
    ("Süre", "~11 s", MUTED, False),
])
_notebox(s,
         "Greedy tüm plakaları yatık istifliyor (anlık en alçak tepe hep yatık "
         "poz); istif sırasını ve dik-plaka fırsatını göremiyor.",
         left=8.2, top=5.0, width=4.6, height=1.7, size=14)
_notes(s,
"'Baseline on bir saniyede 214.5 milimetre. Mantıklı bir yerleşim ama iki körlüğü "
"var: plakaların istif sırasını optimize edemiyor ve kısa plakaları dik dikme "
"fırsatını hiç görmüyor — çünkü o hamle anlık olarak daha yüksek görünüyor.'\n"
"Bu körlük SA'nın değerini kuruyor: bir sonraki slayt.")

# ------------------------------------------------- 6. SA yerleşimi
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Yerleşim 2 — Simulated Annealing (Final)")
img = _RES / "numune_sa" / "nesting3d_numune.png"
if img.exists():
    s.shapes.add_picture(str(img), Inches(0.45), Inches(1.65), width=Inches(7.5))
_caption(s, "Kısa plakalar kenarda dik, uzun plakalar ortada yatık istif", 0.8, 6.9)
_metric_panel(s, [
    ("Z yüksekliği", SA_H, DGREEN, True),
    ("Packing density", SA_D, DGREEN, True),
    ("Kazanç (DBLF'e göre)", GAIN, GREEN, False),
    ("Min boşluk (ölçülen)", CLEAR_MM + " ✓", GREEN, False),
    ("Süre", SA_TIME, MUTED, False),
])
_notes(s,
"'SA, greedy'nin göremediği hamleyi buldu: kısa plakaları kenara dik dikti, "
"uzun plakaları ortada yatık istifledi, küçük parçaları şeritlere dağıttı. "
"Sonuç 181.5 milimetre — yüzde 15.4 kazanç, bugüne kadarki en büyüğü.'\n"
"Üç farklı seed ile koşuldu: 181.5 / 183.0 / 184.5 — sonuç seed'e değil "
"stratejiye bağlı, yani güvenilir. Boşluk her koşuda 1.50 mm ölçüldü.")

# ------------------------------------------------- 7. Karşılaştırma (rakamlar)
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Karşılaştırma — Rakamlar")
rows = [
    ["", "Z yüksekliği", "Density", "Min boşluk", "Not"],
    ["DBLF baseline", DBLF_H, DBLF_D, "—", "greedy, hacim-azalan"],
    ["SA final (2000 iter)", SA_H, SA_D, CLEAR_MM + " ✓", "hibrit poz + sıra araması"],
    ["Önceki 155 mm koşusu", "155.0 mm", "0.118", "0.77 mm ✗", "350 taban, boşluksuz — GEÇERSİZ"],
]
tbl = s.shapes.add_table(4, 5, Inches(0.8), Inches(1.85),
                         Inches(11.7), Inches(2.5)).table
widths = [2.9, 2.0, 1.5, 1.8, 3.5]
for c, w in enumerate(widths):
    tbl.columns[c].width = Inches(w)
for c in range(5):
    cell = tbl.cell(0, c)
    cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
    para = cell.text_frame.paragraphs[0]; run = para.add_run(); run.text = rows[0][c]
    run.font.size = Pt(16); run.font.bold = True; run.font.color.rgb = WHITE
    para.alignment = PP_ALIGN.CENTER
for ri in range(1, 4):
    win = (ri == 2)
    bad = (ri == 3)
    for c in range(5):
        cell = tbl.cell(ri, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = WINBG if win else WHITE
        para = cell.text_frame.paragraphs[0]; run = para.add_run(); run.text = rows[ri][c]
        run.font.size = Pt(15)
        run.font.bold = win or (c == 1)
        run.font.color.rgb = (DGREEN if win else (RED if bad and c in (3, 4)
                                                  else INK))
        para.alignment = PP_ALIGN.CENTER

band = s.shapes.add_shape(1, Inches(0.8), Inches(4.7), Inches(11.7), Inches(1.0))
band.fill.solid(); band.fill.fore_color.rgb = WINBG; band.line.fill.background()
btf = band.text_frame; btf.word_wrap = True
btf.margin_left = Pt(16); btf.margin_top = Pt(10)
bp = btf.paragraphs[0]
br = bp.add_run()
br.text = f"Yeni şartlar altında geçerli en iyi sonuç: {SA_H}  ·  SA kazancı {GAIN}"
br.font.size = Pt(22); br.font.bold = True
br.font.color.rgb = DGREEN
_notebox(s,
         f"Dürüst çerçeve: önceki 155 mm sonucu hem 350 mm tabandaydı hem parça arası "
         f"boşluk içermiyordu — sonradan ölçtük: 0.77 mm'lik İHLAL vardı. Yeni kurallarla "
         f"kıyas geçersiz. {HOCA_HEDEF} sözlü hedefi bu şartlarda ulaşılamadı; en yakın "
         f"geçerli sonuç {SA_H} (kısa plakanın dik boyu ~178 mm yapısal taban oluşturuyor).",
         top=5.95, height=1.3, size=14)
_notes(s,
"'Tabloda üç satır var. İlk ikisi yeni şartlarla bugünkü sonuçlar. Üçüncüsü geçen "
"haftaki 155 — onu bilerek koydum: o koşu daha büyük tabandaydı, boşluk kuralı yoktu "
"ve ölçünce 0.77 milimetrelik ihlal çıktı. Yani 155 ile 181.5 elma-armut.'\n"
"170 SORULURSA: '170'e bu şartlarda ulaşamadık; matematiksel engel şu — en kısa "
"plakanın dik boyu 178.3, yatık istifin tabanı da plaka kalınlıkları toplamından "
"geliyor. 170'in altı ancak taban büyürse veya plaka geometrisi geçişmeye uygun "
"olsaydı mümkündü. İmkânsız demiyorum ama bu konfigürasyonun yapısal sınırındayız.'")

# ------------------------------------------------- 8. Görsel karşılaştırma
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Karşılaştırma — Yerleşimler Yan Yana")
_pic_center_h(s, _RES / "numune_comparison.png", top=1.6, height=5.7)
_notes(s,
"'Solda baseline, sağda SA — aynı parçalar, aynı taban. Yükseklik 214.5'ten "
"181.5'e iniyor, density 0.189'dan 0.229'a çıkıyor.'\n"
"Görseldeki fark: SA kısa plakaları kenara dik dikiyor (sağ görselde dik duran "
"levhalar) ve uzun plakaları ortada derli toplu istifliyor.")

# ------------------------------------------------- 9. Yakınsama
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "SA Yakınsaması — Kazanç Nereden Geldi?")
img = _RES / "numune_sa" / "sa3d_convergence_numune.png"
if img.exists():
    s.shapes.add_picture(str(img), Inches(0.55), Inches(1.8), width=Inches(7.6))
_caption(s, "Best-so-far yükseklik / iterasyon — kesikli çizgi DBLF baseline", 0.8, 6.3)
tb = s.shapes.add_textbox(Inches(8.4), Inches(2.0), Inches(4.4), Inches(4.2))
tf = tb.text_frame; tf.word_wrap = True
for i, (txt, c, b) in enumerate([
    ("214.5 → 181.5 mm", GREEN, True),
    ("Kazancın motoru: kısa plakaları dik çevirme + plaka istif sırası "
     "hamleleri — SA bunları kabul ede ede tavanı söndürüyor.", INK, False),
    ("Üç seed (42/7/13) → 184.5 / 183.0 / 181.5: bant dar, strateji "
     "seed'e dayanıklı.", INK, False),
    ("2000 iterasyon ~93 dk; süre/kalite ayarı --iters ile elimizde.",
     INK, False),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(12)
    r = p.add_run(); r.text = ("" if i == 0 else "•  ") + txt
    r.font.size = Pt(22 if i == 0 else 16); r.font.color.rgb = c; r.font.bold = b
_notes(s,
"'Eğri baseline'dan başlıyor — SA hiçbir zaman baseline'ın üstüne çıkmaz, "
"en iyi çözüm hep saklanır. Kazanç iki hamle ailesinden geliyor: kısa plakaları "
"dik çevirmek ve plaka istif sırasını değiştirmek.'\n"
"Süre sorulursa: bu bir planlama hesabı, gerçek zamanlı değil — gece koşusu "
"normal. Daha kısa sürede yaklaşık sonuç gerekiyorsa iterasyon düşürülür.")

# ------------------------------------------------- 10. KANIT: parça arası boşluk
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Kanıt 1 — Parça Arası Boşluk (İç İçe Geçme YOK)")
_bullets(s, [
    ("İki katmanlı güvence — önce garanti, sonra ölçüm.", INK, True),
    ("Garanti (voxel modelde): parça yüzeyi örnekleme ile KONSERVATİF sarılır "
     "(ince duvar/kabartma kaçmaz), gövde 1 voxel genişletilir, istifte dikey "
     "boşluk hücresi eklenir → çakışma matematiksel olarak imkânsız.", INK, False),
    (f"Ölçüm (orijinal geometride): yerleşim sonrası her yakın parça çifti "
     f"({CLEAR_PAIRS}) gerçek mesh'ler üzerinde tarandı → "
     f"min mesafe {CLEAR_MM} (eşik 1.00 mm).", GREEN, True),
    ("Rapor: results/numune_sa/clearance_numune.json — sayısal, denetlenebilir.",
     INK, False),
    ("Plakaların 'iç içe' görünmesi: çıkıntılar komşu plakanın delik BOŞLUĞUNA "
     "sarkar; yüzeyler hiçbir noktada 1.50 mm'den fazla yaklaşmaz.", INK, False),
])
_notebox(s,
         f"Geçen haftaki 0.77 mm ihlali bu ölçüm aracı yakaladı ve kök neden "
         f"(voxelization'da ince duvarların kaçması) giderildi. Aynı araç artık her "
         f"koşunun standart adımı: boşluk iddiası ölçümsüz rapor edilmiyor.",
         top=5.8, height=1.1, size=14)
_notes(s,
"'İç içe geçme sorusu için iki cevabım var. Birincisi tasarım: voxel modelde iki "
"parçanın aynı hücreyi kullanması imkânsız ve gövdeler bir hücre şişirilerek "
"yerleştiriliyor. İkincisi kanıt: yerleşim bittikten sonra gerçek geometri üzerinde "
"en yakın kırk altı çiftin mesafesini ölçüyorum — minimum 1.50 milimetre.'\n"
"Bu slayt güven slaytı: kendi hatamızı (0.77) kendi aracımızın yakaladığını "
"söylemek güvenilirliği artırır.")

# ------------------------------------------------- 11. KANIT: çözünürlük
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Kanıt 2 — Çözünürlük: Kabartma Yazılar Korunuyor")
_bullets(s, [
    ("Mimari: optimizasyon voxel'de çalışır, çıktı ORİJİNAL geometriden "
     "üretilir (devoxelization) — iki dünya birbirine karışmaz.", INK, True),
    (f"Voxel ({PITCH}) yalnız yerleşim hesabının İÇ aracıdır: {PLATE} taban "
     f"223×223 hücreye bölünür, yükseklik 1.5 mm katmanlara YUKARI yuvarlanır "
     f"(muhafazakâr). Çıktının çözünürlüğü voxel DEĞİL, orijinal mesh'tir — "
     f"'en az 1 mm çözünürlük' şartı çıktıda fazlasıyla aşılır.", INK, True),
    (f"Sayısal kanıt: 48 parçanın orijinal üçgen toplamı {TRI} — "
     f"çıktı STL'sindeki üçgen sayısı {TRI}. Bire bir; tek üçgen kaybolmuyor.",
     GREEN, True),
    ("Kabartma yazılar, delik kenarları, ince detaylar üretim dosyasında "
     "tam çözünürlükte: results/numune_sa/nesting3d_result_numune.stl.", INK, False),
    ("Ekrandaki hızlı önizleme (PNG) sadeleştirilmiş kopya kullanır — "
     "üretime giden STL'yi ETKİLEMEZ.", INK, False),
    ("Geçen haftaki sorun (yazıların silinmesi) görselleştirme sadeleştirmesinin "
     "yanlışlıkla çıktıya uygulanmasıydı — mimari olarak ayrıldı, tekrarı "
     "mümkün değil.", INK, False),
])
_notes(s,
"'Çözünürlük sorusunun cevabı mimaride: voxel yalnız yerleşim hesabı için var — "
"1.5 milimetrelik hücreler, 335'lik tabanda 223'e 223'lük grid. Yerleşim bitince "
"her parçanın ORİJİNAL mesh'i bulunan pozisyona taşınıyor; yani çıktının "
"çözünürlüğünü voxel değil, sizin gönderdiğiniz STL'nin kendisi belirliyor. "
"Kanıtı da basit: kırk sekiz parçanın üçgen toplamı üç milyon dokuz yüz elli bin — "
"çıktı dosyasında da aynı sayı.'\n"
"'En az 1 mm çözünürlük' sorulursa: voxel 1.5 mm ama o yalnız İÇ hesap; "
"raporlanan yükseklik yukarı yuvarlandığı için muhafazakâr, çıktı ise birebir "
"orijinal. İsterse STL'yi yanında açıp bir plakanın yazısına zoom yap — "
"en ikna edici an bu olur.")

# ------------------------------------------------- 12. Çıktılar + kullanım
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Çıktılar ve Kullanıma Hazırlık")
_bullets(s, [
    ("Üretime hazır STL sahnesi: results/numune_sa/nesting3d_result_numune.stl "
     "(orijinal çözünürlük) + yerleşim görselleri + clearance raporu (JSON).",
     GREEN, True),
    ("Parametrik: taban (--plate), voxel (--pitch), boşluk (--margin), "
     "poz stratejisi (--orient), iterasyon (--iters) — makine/kural değişirse "
     "aynı gün yeniden koşulur.", INK, False),
    ("Farklı parça setleri: Numuneler/ klasörüne STL koymak + adet yazmak "
     "yeterli — pipeline parça-bağımsız.", INK, False),
    ("Doğrulama: 95 birim testi · sabit seed · her koşuda otomatik clearance "
     "ölçümü.", INK, False),
    ("Rapor tabloları: results/numune_comparison.md + numune_parts.csv.", INK, False),
])
_notebox(s,
         "Sıradaki adım önerisi: bu motor, konteyner yerleştirme uygulamasının "
         "çekirdeği olmaya hazır — sabit kasa boyutu + çoklu konteyner sarmalayıcı "
         "+ veri girişi arayüzü ile demo kurulabilir.",
         top=5.8, height=1.1)
_notes(s,
"'Sonuç sadece rakam değil — teslim edilebilir paket: üretime girebilecek STL, "
"denetlenebilir boşluk raporu, tekrar koşulabilir parametrik sistem.'\n"
"Uygulama konusu açılırsa: motor hazır, üstüne konteyner kabuğu + arayüz "
"giydirilecek — demo birkaç hafta mesafede.")

# ------------------------------------------------- 13. Özet
s = prs.slides.add_slide(BLANK); _bg(s); _accent_bar(s)
_title(s, "Özet")
_bullets(s, [
    (f"Makine şartlarıyla ({PLATE}, z ≤ {MAX_Z}, boşluk {MARGIN}) 48 parçanın "
     f"tamamı yerleşti.", INK, False),
    (f"Z yüksekliği: DBLF {DBLF_H} → SA {SA_H}  ({GAIN}).", GREEN, True),
    (f"Parça arası ölçülen min boşluk: {CLEAR_MM} — iç içe geçme/çakışma YOK.",
     GREEN, True),
    (f"Çözünürlük: {TRI} üçgenin tamamı çıktıda — kabartma yazılar korunuyor.",
     GREEN, True),
    ("Önceki 155 mm: farklı taban + boşluksuz + ölçülmüş ihlal → geçersiz kıyas; "
     "yeni şartların geçerli en iyisi 181.5 mm.", INK, False),
    ("Tekrar-üretilebilir: 95 test · sabit seed · parametrik CLI · STL hazır.",
     INK, False),
])
qb = s.shapes.add_textbox(Inches(0.8), Inches(6.0), Inches(11.7), Inches(1))
qr = qb.text_frame.paragraphs[0].add_run()
qr.text = ("Teşekkürler — makine tablası ve kural setinizde birebir doğrulama "
           "koşusu yapalım mı?")
qr.font.size = Pt(24); qr.font.color.rgb = MUTED; qr.font.italic = True
_notes(s,
"KAPANIŞ: 'Şartlarınızın üçü de uygulandı ve ikisi ölçümle kanıtlı. Yeni kurallar "
"altında 48 parça 181.5 milimetrede, 1.50 milimetre ölçülmüş boşlukla, yazılar "
"korunarak yerleşiyor. Bir sonraki adım olarak bunu sizin makinenizde birebir "
"doğrulayalım; sonrasında konteyner uygulamasına aynı motorla geçebiliriz.'")

out = _ROOT / "SUNUM_NUMUNE.pptx"
prs.save(str(out))
print(f"Yazildi: {out}  ({len(prs.slides)} slayt)")
